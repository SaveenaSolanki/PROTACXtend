#!/usr/bin/env python3
"""PyMOL structure views PPT."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

H2 = "/storage/saveena/protacpilot/ICM_HMGB2_Hypothesis_Testing/02_H2_ring_modified_ICM_nanobinder"
FIGS = os.path.join(H2, "proof")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)

def add_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_bar(s, l, t, w, h, text, color=DARK_BLUE, size=24):
    shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    tx = s.shapes.add_textbox(Inches(l+0.3), Inches(t+0.05), Inches(w-0.6), Inches(h-0.1))
    p = tx.text_frame.paragraphs[0]; p.text = text; p.font.size = Pt(size); p.font.bold = True; p.font.color.rgb = WHITE

def add_text(s, l, t, w, h, text, size=13, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT):
    tx = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color; p.alignment = align

def add_img(s, path, l, t, w, h=None):
    if os.path.exists(path):
        if h: s.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
        else: s.shapes.add_picture(path, Inches(l), Inches(t), Inches(w))

# SLIDE 1: TITLE
s = add_slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = DARK_BLUE; bg.line.fill.background()
add_text(s, 1, 1.8, 11, 1.5, "PYMOL STRUCTURE VIEWS", 44, True, WHITE, PP_ALIGN.CENTER)
add_text(s, 1, 3.3, 11, 1, "HMGB2 + A1_4COOH | Binding Site | Residue Interactions | The 96° Problem", 20, False, RGBColor(0xD5,0xDB,0xDB), PP_ALIGN.CENTER)
add_text(s, 1, 4.8, 11, 0.6, "Rendered with PyMOL 3.1.0 | ray-traced at 1200x1000, 300 dpi", 14, False, RGBColor(0x85,0xA5,0xC0), PP_ALIGN.CENTER)

# SLIDE 2: Overview
s = add_slide()
add_bar(s, 0, 0, 13.333, 0.9, "HMGB2 + A1_4COOH Overview (Cartoon)", DARK_BLUE)
add_img(s, os.path.join(FIGS, "pymol_hmgb2_a1_overview.png"), 2.0, 1.1, 9.3)
add_text(s, 0.5, 6.6, 12.3, 0.7, 
    "Blue = Box A | Green = Box B | Purple = C-tail | Orange sticks = A1_4COOH | Gold sticks = ICM site residues (78-86) | Lime sticks = CRBN interface (112-128) | Yellow = COOH",
    11, False, DARK_TEXT, PP_ALIGN.CENTER)

# SLIDE 3: Pocket residues
s = add_slide()
add_bar(s, 0, 0, 13.333, 0.9, "Binding Pocket — Key Residues (PyMOL)", DARK_BLUE)
add_img(s, os.path.join(FIGS, "pymol_pocket_residues.png"), 2.0, 1.1, 9.3)
add_text(s, 0.5, 6.6, 12.3, 0.7,
    "Gray surface = HMGB2 pocket (residues 70-95) | Orange = A1_4COOH | Magenta = key residues (LYS85, TYR78, GLY83) | Yellow dashed = salt bridge COOH→LYS85 | Cyan dashed = H-bonds",
    11, False, DARK_TEXT, PP_ALIGN.CENTER)

# SLIDE 4: Surface
s = add_slide()
add_bar(s, 0, 0, 13.333, 0.9, "A1_4COOH in HMGB2 Surface", DARK_BLUE)
add_img(s, os.path.join(FIGS, "pymol_surface_a1.png"), 2.0, 1.1, 9.3)
add_text(s, 0.5, 6.6, 12.3, 0.7,
    "Transparent skyblue = HMGB2 solvent surface | Gold = ICM binding site surface (78-86) | Orange sticks = A1_4COOH | Yellow sphere = COOH | Shows the COOH reaching out of the pocket",
    11, False, DARK_TEXT, PP_ALIGN.CENTER)

# SLIDE 5: Overlay
s = add_slide()
add_bar(s, 0, 0, 13.333, 0.9, "Pose Overlay: ICM (Red) vs A1_4COOH (Blue)", DARK_BLUE)
add_img(s, os.path.join(FIGS, "pymol_overlay_icm_a1.png"), 2.0, 1.1, 9.3)
add_text(s, 0.5, 6.6, 12.3, 0.7,
    "Red = parent ICM (-5.75) | Blue = A1_4COOH (-11.22) | Yellow = COOH | Same binding site; the COOH modification extends a new handle",
    11, False, DARK_TEXT, PP_ALIGN.CENTER)

# SLIDE 6: 96-degree problem
s = add_slide()
add_bar(s, 0, 0, 13.333, 0.9, "The 96° Geometry Problem (PyMOL Surface)", ACCENT_RED)
add_img(s, os.path.join(FIGS, "pymol_96deg_problem.png"), 2.0, 1.1, 9.3)
add_text(s, 0.5, 6.6, 12.3, 0.7,
    "Gold surface = ICM binding site (78-86) | Green surface = CRBN interface (112-128) | Red dashed = distance between sites | Explains why the PROTAC cannot reach CRBN",
    11, False, DARK_TEXT, PP_ALIGN.CENTER)

# SLIDE 7: Lysines
s = add_slide()
add_bar(s, 0, 0, 13.333, 0.9, "HMGB2 Lysine Landscape (Purple Spheres)", DARK_BLUE)
add_img(s, os.path.join(FIGS, "pymol_lysines_hmgb2.png"), 2.0, 1.1, 9.3)
add_text(s, 0.5, 6.6, 12.3, 0.7,
    "Purple spheres = 40 lysine sidechains | Orange = A1_4COOH | All lysines accessible for ubiquitination — key for S1 (glue) and S5 (degron) strategies",
    11, False, DARK_TEXT, PP_ALIGN.CENTER)

# SLIDE 8: Gallery
s = add_slide()
add_bar(s, 0, 0, 13.333, 0.9, "All PyMOL Renderings", DARK_BLUE)
views = [
    ("pymol_hmgb2_a1_overview.png", "Overview"),
    ("pymol_pocket_residues.png", "Pocket Residues"),
    ("pymol_surface_a1.png", "Surface"),
    ("pymol_overlay_icm_a1.png", "ICM vs A1"),
    ("pymol_96deg_problem.png", "96° Problem"),
    ("pymol_lysines_hmgb2.png", "Lysines"),
]
for i, (fname, label) in enumerate(views):
    row, col = i // 3, i % 3
    x = 0.8 + col * 4.2; y = 1.1 + row * 3.0
    img_path = os.path.join(FIGS, fname)
    if os.path.exists(img_path):
        try: s.shapes.add_picture(img_path, Inches(x), Inches(y), Inches(3.8), Inches(2.4))
        except: pass
    add_text(s, x, y + 2.5, 3.8, 0.3, label, 10, True, DARK_TEXT, PP_ALIGN.CENTER)

ppt_path = os.path.join(H2, "H2_PyMOL_Structure_Views.pptx")
prs.save(ppt_path)
print(f"✅ Saved: {ppt_path} ({len(prs.slides)} slides)")
