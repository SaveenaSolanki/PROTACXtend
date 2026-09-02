#!/usr/bin/env python3
"""
Structure Views PPT - PDB-based visualizations of HMGB2 binding sites and poses.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

H2 = "/storage/saveena/protacpilot/ICM_HMGB2_Hypothesis_Testing/02_H2_ring_modified_ICM_nanobinder"
FIGS = os.path.join(H2, "proof")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)

def add_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_bar(s, l, t, w, h, text, color=DARK_BLUE, size=24):
    shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    tx = s.shapes.add_textbox(Inches(l+0.3), Inches(t+0.05), Inches(w-0.6), Inches(h-0.1))
    p = tx.text_frame.paragraphs[0]; p.text = text; p.font.size = Pt(size); p.font.bold = True; p.font.color.rgb = WHITE

def add_text(s, l, t, w, h, text, size=13, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT):
    tx = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color; p.alignment = align

def add_img(s, path, l, t, w, h=None):
    if os.path.exists(path):
        if h: s.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
        else: s.shapes.add_picture(path, Inches(l), Inches(t), Inches(w))

# ======================================================================
# SLIDE 1: TITLE
# ======================================================================
s = add_slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = DARK_BLUE; bg.line.fill.background()
add_text(s, 1, 1.8, 11, 1.5, "HMGB2 STRUCTURE VIEWS", 42, True, WHITE, PP_ALIGN.CENTER)
add_text(s, 1, 3.3, 11, 1, "PDB-Based 3D Visualizations of Binding Sites, Poses, and Key Residues", 20, False, RGBColor(0xD5,0xDB,0xDB), PP_ALIGN.CENTER)
add_text(s, 1, 4.8, 11, 0.6, "Source PDB: hmgb2_fixed_minim.pdb (AlphaFold + P4ward minimization)", 14, False, RGBColor(0x85,0xA5,0xC0), PP_ALIGN.CENTER)

# ======================================================================
# SLIDE 2: HMGB2 3D structure with sites
# ======================================================================
s = add_slide()
add_bar(s, 0, 0, 13.333, 0.9, "HMGB2 3D Structure: A1_4COOH + Binding Sites", DARK_BLUE)
add_img(s, os.path.join(FIGS, "pdb_hmgb2_sites_3d.png"), 2.5, 1.1, 8.3)
add_text(s, 0.5, 6.7, 12.3, 0.6, 
    "Blue = Box A | Green = Box B | Purple = C-tail | Orange = A1_4COOH | Gold sphere = ICM site (78-86) | Green sphere = CRBN interface (112-128) | Red dashed = 96° apart",
    11, False, DARK_TEXT, PP_ALIGN.CENTER)

# ======================================================================
# SLIDE 3: Binding pocket close-up
# ======================================================================
s = add_slide()
add_bar(s, 0, 0, 13.333, 0.9, "A1_4COOH Binding Pocket — 3D Close-up", DARK_BLUE)
add_img(s, os.path.join(FIGS, "pdb_pocket_3d.png"), 2.5, 1.1, 8.3)
add_text(s, 0.5, 6.7, 12.3, 0.6,
    "Gray/colored spheres = protein atoms in pocket | Orange = A1_4COOH | Gold diamond = COOH | Pink = key residues (TYR78 H-bond, LYS85 salt bridge, LYS82/86 basic)",
    11, False, DARK_TEXT, PP_ALIGN.CENTER)

# ======================================================================
# SLIDE 4: Pose overlay ICM vs A1_4COOH
# ======================================================================
s = add_slide()
add_bar(s, 0, 0, 13.333, 0.9, "Pose Overlay: ICM vs A1_4COOH in HMGB2", DARK_BLUE)
add_img(s, os.path.join(FIGS, "pdb_pose_overlay_3d.png"), 2.5, 1.1, 8.3)
add_text(s, 0.5, 6.7, 12.3, 0.6,
    "Red = parent ICM (-5.75) | Blue = A1_4COOH (-11.22) | Gold = COOH extension | Both bind at the same site (78-86); A1_4COOH adds the COOH handle",
    11, False, DARK_TEXT, PP_ALIGN.CENTER)

# ======================================================================
# SLIDE 5: Surface views (2 projections)
# ======================================================================
s = add_slide()
add_bar(s, 0, 0, 13.333, 0.9, "HMGB2 Surface Views with Ligand", DARK_BLUE)
add_img(s, os.path.join(FIGS, "pdb_surface_views.png"), 0.5, 1.1, 12.3)
add_text(s, 0.5, 6.7, 12.3, 0.6,
    "Left = XZ projection | Right = YZ projection | Cyan = protein surface | Orange = A1_4COOH | Gold = ICM site | Green = CRBN interface",
    11, False, DARK_TEXT, PP_ALIGN.CENTER)

# ======================================================================
# SLIDE 6: Site residue map
# ======================================================================
s = add_slide()
add_bar(s, 0, 0, 13.333, 0.9, "HMGB2 Domain Map with All Key Sites", DARK_BLUE)
add_img(s, os.path.join(FIGS, "pdb_site_residue_map.png"), 0.5, 1.3, 12.3)
add_text(s, 0.5, 6.8, 12.3, 0.6,
    "Blue = Box A (1-79) | Green = Box B (95-163) | Purple = C-tail (164-209) | Gold = ICM site (78-86) | Green = CRBN interface (112-128) | Purple triangles = 40 lysines | Stars = dTAG insertion points",
    11, False, DARK_TEXT, PP_ALIGN.CENTER)

# ======================================================================
# SLIDE 7: All structure views gallery
# ======================================================================
s = add_slide()
add_bar(s, 0, 0, 13.333, 0.9, "Complete Structure View Gallery", DARK_BLUE)

views = [
    ("pdb_hmgb2_sites_3d.png", "Sites 3D"),
    ("pdb_pocket_3d.png", "Pocket 3D"),
    ("pdb_pose_overlay_3d.png", "Pose Overlay"),
    ("pdb_surface_views.png", "Surface Views"),
    ("pdb_site_residue_map.png", "Residue Map"),
    ("residue_level_interactions.png", "Interactions"),
    ("binding_pose_detailed.png", "Binding Pose"),
    ("boltz_interface.png", "Boltz Interface"),
]
for i, (fname, label) in enumerate(views):
    row, col = i // 4, i % 4
    x = 0.4 + col * 3.2; y = 1.1 + row * 2.9
    img_path = os.path.join(FIGS, fname)
    if os.path.exists(img_path):
        try: s.shapes.add_picture(img_path, Inches(x), Inches(y), Inches(2.9), Inches(2.2))
        except: pass
    add_text(s, x, y + 2.3, 2.9, 0.3, label, 9, True, DARK_TEXT, PP_ALIGN.CENTER)

# ======================================================================
# SAVE
# ======================================================================
ppt_path = os.path.join(H2, "H2_Structure_Views_Presentation.pptx")
prs.save(ppt_path)
print(f"✅ Saved: {ppt_path} ({len(prs.slides)} slides)")
