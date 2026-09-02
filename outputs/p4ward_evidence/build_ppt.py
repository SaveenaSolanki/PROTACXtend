#!/usr/bin/env python3
"""
Build complete 11-slide HMGB2 PROTAC meeting presentation.
Output: HMGB2_PROTAC_Meeting.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = "/storage/saveena/protacpilot/outputs/p4ward_evidence"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Color scheme
DARK_BLUE = RGBColor(0x1F, 0x38, 0x64)
MED_BLUE = RGBColor(0x2F, 0x54, 0x96)
LIGHT_BLUE = RGBColor(0x44, 0x72, 0xC4)
ACCENT_RED = RGBColor(0xC0, 0x00, 0x00)
ACCENT_ORANGE = RGBColor(0xED, 0x7D, 0x31)
ACCENT_GREEN = RGBColor(0x54, 0x8C, 0x2F)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


def add_background(slide, color=LIGHT_GRAY):
    """Add a solid background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, line_color=None):
    """Add a colored rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_text(slide, left, top, width, height, items, font_size=16,
                    color=BLACK, font_name='Calibri', spacing=Pt(6)):
    """Add bulleted text items."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        p.level = 0
    return txBox


def add_image(slide, path, left, top, width=None, height=None):
    """Add an image to the slide."""
    if not os.path.exists(path):
        # Placeholder
        shape = add_shape(slide, left, top, width or Inches(4), height or Inches(3), 
                          RGBColor(0xE0, 0xE0, 0xE0))
        add_textbox(slide, left + Inches(0.5), top + Inches(1), 
                    (width or Inches(4)) - Inches(1), Inches(1),
                    f"[Image: {os.path.basename(path)}]", 
                    font_size=11, color=RGBColor(0x88, 0x88, 0x88))
        return None
    return slide.shapes.add_picture(path, left, top, width=width, height=height)


def add_title_bar(slide, title_text, subtitle_text=None):
    """Add a consistent title bar at the top."""
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), DARK_BLUE)
    add_textbox(slide, Inches(0.5), Inches(0.1), Inches(12), Inches(0.6),
                title_text, font_size=28, bold=True, color=WHITE)
    if subtitle_text:
        add_textbox(slide, Inches(0.5), Inches(0.6), Inches(12), Inches(0.4),
                    subtitle_text, font_size=14, color=RGBColor(0xBB, 0xCC, 0xDD))


def add_footer(slide, slide_num, total=11):
    """Add slide number footer."""
    add_textbox(slide, Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.3),
                f"{slide_num}/{total}", font_size=10, color=RGBColor(0x99, 0x99, 0x99),
                alignment=PP_ALIGN.RIGHT)


# ======================================================================
# SLIDE 1 — Objective
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_background(slide)
add_title_bar(slide, "Why Did HMGB2 PROTACs Fail?", "Target biology vs. PROTAC geometry")

# Left panel: HMGB2 biology
add_shape(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.5), WHITE)
add_textbox(slide, Inches(0.7), Inches(1.4), Inches(5.4), Inches(0.4),
            "HMGB2 — The Target", font_size=20, bold=True, color=DARK_BLUE)

hmgb2_items = [
    "Nuclear chromatin-binding protein (Box A + Box B domains)",
    "209 amino acids, 24 kDa, pI ~9.5 (highly basic)",
    "20+ surface-accessible lysines — favorable for ubiquitination",
    "Long, flexible, acidic C-terminal tail (186–209 aa)",
    "Rapid dynamic exchange on/off chromatin (FRAP t½ ~seconds)",
    "Involved in: inflammation, cancer, chromatin remodeling",
    "Unregulated in multiple cancers (prognostic marker)",
]
add_bullet_text(slide, Inches(0.7), Inches(1.9), Inches(5.4), Inches(4.5),
                hmgb2_items, font_size=14, spacing=Pt(4))

# Right panel: The question
add_shape(slide, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.5), WHITE)
add_textbox(slide, Inches(7.0), Inches(1.4), Inches(5.6), Inches(0.4),
            "The PROTAC Design Question", font_size=20, bold=True, color=DARK_BLUE)

question_items = [
    "Warhead: Inflachromene (ICM) — known HMGB1/2 binder (Lee et al., 2014)",
    "E3 ligases: VHL (AHPC) and CRBN (Thalidomide)",
    "Linkers: C4, C6, C8 alkyl (3 lengths × 2 E3 = 6 PROTACs)",
    "",
    "RESULT: All 6 PROTACs showed NO degradation",
    "",
    "Core questions:",
    "  1. Is the ternary complex geometrically possible?",
    "  2. Is HMGB2 intrinsically non-degradable?",
    "  3. Which design parameter caused the failure?",
]
add_bullet_text(slide, Inches(7.0), Inches(1.9), Inches(5.6), Inches(4.5),
                question_items, font_size=14, spacing=Pt(3))

# Highlight box
add_shape(slide, Inches(7.0), Inches(5.5), Inches(5.4), Inches(1.0), 
          RGBColor(0xFF, 0xF2, 0xCC), RGBColor(0xC0, 0x00, 0x00))
add_textbox(slide, Inches(7.2), Inches(5.6), Inches(5.0), Inches(0.8),
            "This talk: Systematic computational diagnosis\n→ Root cause identified → Rational redesign plan",
            font_size=13, bold=True, color=ACCENT_RED)

add_footer(slide, 1)

# ======================================================================
# SLIDE 2 — Pipeline
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "Computational Pipeline", "HMGB2 modeling workflow — from structure to redesign")

add_image(slide, f"{OUT}/plot08_pipeline.png",
          Inches(1.5), Inches(1.4), width=Inches(10), height=Inches(2.5))

# Detailed explanation
add_shape(slide, Inches(0.5), Inches(4.0), Inches(12.3), Inches(3.0), WHITE)
add_textbox(slide, Inches(0.7), Inches(4.1), Inches(11.9), Inches(0.4),
            "Pipeline Steps", font_size=18, bold=True, color=DARK_BLUE)

steps = [
    "1. HMGB2 + CRBN structure preparation (fix, minimize, protonate)  →  PDB prep",
    "2. MegaDock: 3600 orientations of CRBN around HMGB2  →  exhaustive sampling",
    "3. PROTAC linker distance filter (auto-cutoff: 0.74 Å)  →  checks geometric feasibility",
    "4. Ubiquitination filter (CRBN: 16 Å cutoff)  →  checks lysine-to-E2~Ub distance",
    "5. Vina docking: 12 warheads against HMGB2  →  score + rank alternatives",
    "6. Root-cause diagnosis + redesign proposal  →  next cycle",
]
add_bullet_text(slide, Inches(0.7), Inches(4.6), Inches(11.9), Inches(2.2),
                steps, font_size=13, spacing=Pt(2))

add_footer(slide, 2)

# ======================================================================
# SLIDE 3 — Original P4ward Result
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "P4ward Ternary Complex Modeling — Primary Result",
              "3600 orientations sampled → 0 viable ternary complexes")

# Main figure
add_image(slide, f"{OUT}/plot01_filtering_result.png",
          Inches(0.5), Inches(1.3), width=Inches(5.5), height=Inches(4.5))

# Key stats panel
add_shape(slide, Inches(6.5), Inches(1.3), Inches(6.3), Inches(5.5), WHITE,
          RGBColor(0xCC, 0x00, 0x00))
add_textbox(slide, Inches(6.7), Inches(1.4), Inches(5.9), Inches(0.4),
            "KEY RESULT", font_size=22, bold=True, color=ACCENT_RED)

stats = [
    "Total MegaDock orientations:   3,600",
    "Passed linker distance filter:         0",
    "Passed ubiquitination filter:          0",
    "",
    "Closest exit-vector gap:    10.83 Å",
    "Linker max conformational span:    0.74 Å",
    "Gap vs. linker:                     14.6× too large",
    "",
    "Verification: P4ward log entry",
    '  "There are no poses which satisfy the',
    '   ligand distance filtering criteria.',
    '   Exiting now."',
    "",
    "Method: Sharma et al., Sci Rep 15:21502, 2025",
    "(PRosettaC/P4ward benchmark — best-in-class",
    " for PROTAC ternary complex prediction)",
]
add_bullet_text(slide, Inches(6.7), Inches(1.9), Inches(5.9), Inches(4.5),
                stats, font_size=14, spacing=Pt(1), font_name='Consolas')

# Bottom summary bar
add_shape(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.7), ACCENT_RED)
add_textbox(slide, Inches(0.7), Inches(6.3), Inches(11.9), Inches(0.5),
            "DEFINITIVE RESULT: The C4-equivalent linker (CCCOCCC, 0.74 Å max span) cannot bridge HMGB2 and CRBN → no ternary complex possible",
            font_size=15, bold=True, color=WHITE)

add_footer(slide, 3)

# ======================================================================
# SLIDE 4 — Failed Docked Pose Overview
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "Closest Failed MegaDock/P4ward Pose (#2615)",
              "Best of 3600 orientations — still no viable ternary complex")

# Main PDB image
add_image(slide, f"{OUT}/pymol_pose_2615_labeled.png",
          Inches(0.3), Inches(1.2), width=Inches(12.7), height=Inches(4.5))

# Annotation below
add_shape(slide, Inches(0.5), Inches(5.9), Inches(12.3), Inches(1.2), WHITE,
          RGBColor(0xCC, 0x00, 0x00))

annotations = [
    "GREEN: HMGB2 (receptor, fixed position)  |  PINK: CRBN (ligase, rotated by MegaDock)",
    "Exit vectors (red spheres on ICM, purple spheres on thalidomide) are still 10.83 Å apart",
    "CCCOCCC linker can span at most 0.74 Å → this orientation cannot form a productive ternary complex",
]
add_bullet_text(slide, Inches(0.7), Inches(6.0), Inches(11.9), Inches(1.0),
                annotations, font_size=13, spacing=Pt(2), color=ACCENT_RED)

add_footer(slide, 4)

# ======================================================================
# SLIDE 5 — Failed Pose Zoom
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "Exit Vector Gap — Zoom View",
              "ICM warhead exit vectors vs. thalidomide exit vectors")

add_image(slide, f"{OUT}/pymol_gap_zoom.png",
          Inches(0.3), Inches(1.2), width=Inches(8.0), height=Inches(5.5))

# Side panel with measurements
add_shape(slide, Inches(8.8), Inches(1.2), Inches(4.0), Inches(5.5), WHITE,
          RGBColor(0xCC, 0x00, 0x00))

items = [
    "GAP MEASUREMENT",
    "",
    "ICM OH (exit vector 1)",
    "  ↕  10.83 Å",
    "Thalidomide C (exit vec.)",
    "",
    "vs.",
    "",
    "Linker max span: 0.74 Å",
    "",
    "Gap is 14.6× larger",
    "than linker can span",
    "",
    "→ Geometry impossible",
    "→ All 3600 poses failed",
]
add_bullet_text(slide, Inches(9.0), Inches(1.4), Inches(3.6), Inches(5.0),
                items, font_size=14, spacing=Pt(2), color=ACCENT_RED)

add_footer(slide, 5)

# ======================================================================
# SLIDE 6 — Distance Distribution
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "Distance Distribution — All 3600 MegaDock Poses",
              "Exit-vector gaps range from 10.83 Å to 176.04 Å — none pass the filter")

# Main histogram
add_image(slide, f"{OUT}/plot02_distance_histogram.png",
          Inches(0.3), Inches(1.2), width=Inches(8.0), height=Inches(4.8))

# Closest 10 bar chart
add_image(slide, f"{OUT}/plot09_closest_10.png",
          Inches(8.5), Inches(1.2), width=Inches(4.5), height=Inches(3.2))

# Summary table
add_shape(slide, Inches(8.5), Inches(4.6), Inches(4.5), Inches(2.0), WHITE)
table_text = [
    "Gap Range     Count    Outcome",
    "───────────────────────────────",
    "10–20 Å        36      FAILED",
    "20–30 Å        67      FAILED",
    "30–50 Å       315      FAILED",
    "50–100 Å    1,542     FAILED",
    "100–176 Å  1,640     FAILED",
    "───────────────────────────────",
    "TOTAL       3,600     0 PASSED",
]
add_bullet_text(slide, Inches(8.7), Inches(4.7), Inches(4.1), Inches(1.8),
                table_text, font_size=11, spacing=Pt(0), font_name='Consolas', color=BLACK)

# Bottom note
add_textbox(slide, Inches(0.5), Inches(6.4), Inches(12), Inches(0.5),
            "Even the best 36 poses (top 1%) have exit-vector gaps of 10–20 Å — all far beyond the linker's 0.74 Å limit.",
            font_size=14, bold=True, color=ACCENT_RED)

add_footer(slide, 6)

# ======================================================================
# SLIDE 7 — Vina Warhead Docking
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "Warhead Selection — Vina Docking Screen (12 Warheads vs. HMGB2)",
              "Identifying stronger alternatives to Inflachromene")

add_image(slide, f"{OUT}/plot04_vina_docking.png",
          Inches(0.3), Inches(1.2), width=Inches(7.0), height=Inches(4.5))

# Interpretation panel
add_shape(slide, Inches(7.8), Inches(1.2), Inches(5.0), Inches(5.5), WHITE)

add_textbox(slide, Inches(8.0), Inches(1.3), Inches(4.6), Inches(0.4),
            "Key Findings", font_size=18, bold=True, color=DARK_BLUE)

findings = [
    "Hoechst 33258: −6.49 kcal/mol (best)",
    "  → DNA minor groove binder",
    "  → Engages HMGB2's DNA-binding surface",
    "",
    "PDS (Pyridostatin): −5.87 kcal/mol",
    "  → G-quadruplex ligand",
    "  → Alternative binding mode",
    "",
    "Inflachromene: −5.79 kcal/mol (7th)",
    "  → Only modest affinity",
    "  → Binding site not resolved structurally",
    "",
    "RECOMMENDATION:",
    "Replace ICM with Hoechst 33258 or PDS",
    "for stronger, validated HMGB2 engagement",
]
add_bullet_text(slide, Inches(8.0), Inches(1.8), Inches(4.6), Inches(4.5),
                findings, font_size=12, spacing=Pt(1))

add_footer(slide, 7)

# ======================================================================
# SLIDE 8 — E3 Ligase Check
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "E3 Ligase Selection — VHL vs. CRBN for a Nuclear Target",
              "Subcellular localization determines E3 accessibility to HMGB2")

add_image(slide, f"{OUT}/plot06_e3_comparison.png",
          Inches(0.5), Inches(1.3), width=Inches(7.5), height=Inches(4.0))

# Right panel: recommendation
add_shape(slide, Inches(8.5), Inches(1.3), Inches(4.3), Inches(5.5), WHITE)
add_textbox(slide, Inches(8.7), Inches(1.4), Inches(3.9), Inches(0.4),
            "Decision for HMGB2", font_size=18, bold=True, color=DARK_BLUE)

decision_items = [
    "HMGB2 is NUCLEAR (chromatin-bound)",
    "  → Requires E3 with nuclear access",
    "",
    "CRBN (cereblon):",
    "  • Imported into nucleus via KPNB1",
    "  • Proven nuclear substrate degradation",
    "    (IKZF1, IKZF3, GSPT1)",
    "  • Intramolecular folding improves",
    "    permeability (Poongavanam 2025)",
    "  ✅ RECOMMENDED",
    "",
    "VHL (pVHL):",
    "  • Predominantly cytoplasmic",
    "  • Requires target to shuttle to cytoplasm",
    "  • Suboptimal for nuclear HMGB2",
    "  ❌ Not recommended as primary E3",
    "",
    "Preferred ligand: Pomalidomide > Lenalidomide",
    "  (better CRBN binding, different exit vector)",
]
add_bullet_text(slide, Inches(8.7), Inches(1.9), Inches(3.9), Inches(4.5),
                decision_items, font_size=11, spacing=Pt(1))

add_footer(slide, 8)

# ======================================================================
# SLIDE 9 — Root-Cause Summary
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "Root-Cause Diagnosis — Ranked by Probability",
              "Why the current 6 PROTAC designs failed")

add_image(slide, f"{OUT}/plot05_root_causes.png",
          Inches(0.3), Inches(1.2), width=Inches(7.5), height=Inches(4.5))

# Action plan panel
add_shape(slide, Inches(8.3), Inches(1.2), Inches(4.7), Inches(5.5), WHITE)
add_textbox(slide, Inches(8.5), Inches(1.3), Inches(4.3), Inches(0.4),
            "Action Plan", font_size=18, bold=True, color=DARK_BLUE)

actions = [
    "🔴 HIGH (fix immediately):",
    "  1. Longer linkers (C10–C14, PEG-based)",
    "  2. Resolve ICM binding mode (dock + MD)",
    "  3. Design for ternary cooperativity",
    "",
    "🟡 MEDIUM (address in parallel):",
    "  4. Measure cellular permeability",
    "  5. Check E3 expression in target line",
    "  6. Test wider concentration range",
    "",
    "🟢 LOW (rule out last):",
    "  7. HMGB2 chromatin access",
    "  8. HMGB2 intrinsic degradability",
    "",
    "Note: HMGB2 has 20+ surface lysines",
    "and is small (24 kDa) → likely degradable",
    "if the ternary complex can be formed.",
]
add_bullet_text(slide, Inches(8.5), Inches(1.8), Inches(4.3), Inches(4.5),
                actions, font_size=11, spacing=Pt(1))

add_footer(slide, 9)

# ======================================================================
# SLIDE 10 — Proposed Redesign
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "PROTAC Redesign — Next Cycle", "Evidence-based design matrix: current → proposed fix")

add_image(slide, f"{OUT}/plot07_redesign_matrix.png",
          Inches(0.5), Inches(1.3), width=Inches(12.3), height=Inches(3.0))

# Specific linker proposals
add_shape(slide, Inches(0.5), Inches(4.5), Inches(6.0), Inches(2.5), WHITE)
add_textbox(slide, Inches(0.7), Inches(4.6), Inches(5.6), Inches(0.3),
            "Linker Design Proposals", font_size=16, bold=True, color=DARK_BLUE)

linker_items = [
    "C10-PEG₄-amide:  ~12 Å, PEG improves solubility",
    "C12-alkyl-triazole:  ~13 Å, semi-rigid + H-bond acceptor",
    "C14-PEG₅:  ~16 Å, maximum span test",
    "C8-PEG₃-piperidine:  conformational pre-organization",
]
add_bullet_text(slide, Inches(0.7), Inches(5.0), Inches(5.6), Inches(1.8),
                linker_items, font_size=13, spacing=Pt(3))

# Validation plan
add_shape(slide, Inches(7.0), Inches(4.5), Inches(6.0), Inches(2.5), WHITE)
add_textbox(slide, Inches(7.2), Inches(4.6), Inches(5.6), Inches(0.3),
            "Validation Gates Before Synthesis", font_size=16, bold=True, color=DARK_BLUE)

validation_items = [
    "1. Re-run P4ward with new designs (C10–C14 + Pomalidomide)",
    "2. MD simulation of top ternary complex (100+ ns)",
    "3. Lysine-to-E2~Ub distance check for viable poses",
    "4. Permeability prediction (bRo5 chameleonic behavior)",
    "5. Only if ≥50% poses pass → proceed to synthesis",
]
add_bullet_text(slide, Inches(7.2), Inches(5.0), Inches(5.6), Inches(1.8),
                validation_items, font_size=12, spacing=Pt(2))

add_footer(slide, 10)

# ======================================================================
# SLIDE 11 — Conclusion
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "Conclusion", "Current design failed by geometry — HMGB2 degradability is NOT ruled out")

# Main conclusion box
add_shape(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.5), WHITE,
          RGBColor(0x2F, 0x54, 0x96))

conclusions = [
    "1. P4ward result is definitive: 0/3600 poses passed the linker distance filter.",
    "   → The C4-equivalent linker (0.74 Å) cannot span the HMGB2–CRBN gap (min 10.83 Å).",
    "",
    "2. HMGB2 is likely degradable — 20+ surface lysines, flexible C-tail, small size.",
    "   → The target is not the problem. The PROTAC linker is the problem.",
    "",
    "3. Clear redesign path exists: longer linkers (C10–C14), CRBN-based (pomalidomide),",
    "   alternative warheads (Hoechst 33258 / PDS), re-run P4ward before synthesis.",
]
add_bullet_text(slide, Inches(0.7), Inches(1.6), Inches(11.9), Inches(2.2),
                conclusions, font_size=14, spacing=Pt(2), color=BLACK)

# Next steps
add_shape(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.5), WHITE)
add_textbox(slide, Inches(0.7), Inches(4.4), Inches(11.9), Inches(0.4),
            "Immediate Next Steps", font_size=18, bold=True, color=DARK_BLUE)

next_steps = [
    "1. Extract reconstructed P4ward poses (hmgb2_pose_2615.pdb, crbn_pose_2615.pdb) for team inspection",
    "2. Design C10–C14 PEG-linker PROTACs with pomalidomide (CRBN) and Hoechst 33258/PDS warheads",
    "3. Re-run P4ward with new designs — target: >50% of poses passing the linker filter",
    "4. MD simulation of top candidate(s) — 100+ ns to verify ternary complex stability",
    "5. If computational validation passes → synthesize top 3–5 PROTACs, test in cellular degradation assay",
    "",
    "Evidence package location: outputs/p4ward_evidence/ (all PDBs, MOL2s, logs, plots, PyMOL scripts)",
]
add_bullet_text(slide, Inches(0.7), Inches(4.9), Inches(11.9), Inches(1.8),
                next_steps, font_size=13, spacing=Pt(2))

# Bottom bar
add_shape(slide, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4), 
          RGBColor(0x1F, 0x38, 0x64))
add_textbox(slide, Inches(0.7), Inches(6.85), Inches(11.9), Inches(0.3),
            "Failure is diagnostic, not terminal — the geometry failure tells us exactly what to fix.",
            font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_footer(slide, 11)

# ======================================================================
# SAVE
# ======================================================================
output_path = os.path.join(OUT, "HMGB2_PROTAC_Meeting.pptx")
prs.save(output_path)
print(f"✓ Presentation saved: {output_path}")
print(f"  {len(prs.slides)} slides")
print(f"  Size: {os.path.getsize(output_path) / 1024:.0f} KB")
